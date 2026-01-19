#!/usr/bin/env python3
"""
Create a PowerPoint presentation from the Menopause and the Heart research paper
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(31, 56, 100)  # Dark blue
SUBTITLE_COLOR = RGBColor(68, 114, 196)  # Medium blue
TEXT_COLOR = RGBColor(64, 64, 64)  # Dark gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = SUBTITLE_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title="", right_title=""):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    if left_title:
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_after = Pt(12)

    for i, point in enumerate(left_content):
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    if right_title:
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_after = Pt(12)

    for i, point in enumerate(right_content):
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    return slide

# Slide 1: Title
add_title_slide(prs, "Menopause and the Heart",
                "Chileshe Nkonde-Price, MD & Jeffrey R. Bender, MD")

# Slide 2: Introduction
add_content_slide(prs, "Introduction", [
    "Cardiovascular disease (CVD) is the leading cause of death in US women",
    "Premenopausal women are relatively protected against CVD compared to men",
    "The gender gap narrows at menopause, with CVD incidence increasing sharply",
    "This led to the belief that estrogens are cardioprotective",
    "However, randomized clinical trials have not consistently supported this belief"
])

# Slide 3: Key Points
add_content_slide(prs, "Key Points", [
    "HRT is NOT currently recommended solely to prevent future heart attacks",
    "For women with life-disrupting symptoms: topical estrogens first, then hormone patches at lowest effective dose",
    "Treatment should be maintained for the shortest duration possible",
    "For women a decade+ after menopause without symptoms, HRT should be discontinued",
    "Whether HRT increases risk for heart attacks, strokes, and breast cancer remains under investigation"
])

# Slide 4: Observational Studies
add_content_slide(prs, "Observational Studies", [
    "Nurses Health Study (NHS) - largest observational study",
    "Started in 1976 with 122,000 nurses aged 30-55 years",
    "Surveyed every 2 years over 4-year follow-up (93% participation)",
    "Findings:",
    "  • Women lacking endogenous estrogen had greater CVD risk",
    "  • HRT reduced CVD incidence in postmenopausal women",
    "Criticism: Selection bias - HRT users possibly healthier than non-users"
])

# Slide 5: Randomized Clinical Trials - Overview
add_content_slide(prs, "Randomized Clinical Trials", [
    "Designed to address selection bias in observational studies",
    "PEPI Trial (1995):",
    "  • 875 women aged 45-64 years",
    "  • HRT improved CVD risk markers (↓LDL, ↓fibrinogen, ↑HDL)",
    "Two most influential trials:",
    "  • HERS - Secondary prevention trial",
    "  • WHI - Primary prevention trial"
])

# Slide 6: HERS Study
add_content_slide(prs, "HERS: Secondary Prevention Trial", [
    "2,763 postmenopausal women with established CAD",
    "Age range: 44-79 years (mean 67 years)",
    "Randomized to estrogen/progestin therapy vs. placebo",
    "Results:",
    "  • NO significant reduction in MI or CAD-related death",
    "  • MORE CAD events in first year of HRT",
    "  • FEWER events in years 4 and beyond",
    "Conclusion: Time-dependent effect - early harm, possible late benefit"
])

# Slide 7: WHI Study
add_content_slide(prs, "WHI: Primary Prevention Trial", [
    "27,000 healthy postmenopausal women without CVD",
    "Age range: 50-79 years (mean 63 years, 12.5 years post-menopause)",
    "Planned for 8.5 years, stopped at 5 years due to adverse outcomes",
    "Results:",
    "  • 0.07% INCREASE in cardiovascular events with HRT",
    "  • Similar time trend: more events year 1, fewer in years 4+",
    "Major Impact: Led to widespread avoidance of HRT for CVD prevention"
])

# Slide 8: The Timing Hypothesis
add_content_slide(prs, "Sorting Through the Conundrum", [
    "WHI average enrollment age: 63 years (11-12 years older than typical HRT initiation)",
    "Secondary analysis by age groups (50-59, 60-69, 70-79):",
    "  • Youngest group: LOWEST CVD risk",
    "  • Oldest group: HIGHEST risk",
    "  • HRT appeared PROTECTIVE in youngest group",
    "Meta-analysis of 39,000+ women in 23 trials:",
    "  • HRT reduces CHD risk in women <60 years",
    "  • No benefit in older women",
    "Key Finding: TIMING of intervention is critical"
])

# Slide 9: Recent Trials
add_two_column_slide(prs, "Recent Clinical Trials",
    [
        "KEEPS (2014)",
        "720 women, 42-58 years",
        "Within 36 months of last menses",
        "4-year follow-up",
        "Results:",
        "• No difference in CIMT or CAC",
        "• Improved lipid profiles",
        "• Better insulin sensitivity"
    ],
    [
        "ELITE (2014)",
        "643 women in 2 groups:",
        "  • <6 years post-menopause",
        "  • ≥10 years post-menopause",
        "5-year treatment duration",
        "Results:",
        "• ↓ vascular disease progression",
        "• ONLY if started within 6 years",
        "• Reinforces timing hypothesis"
    ],
    "KEEPS", "ELITE"
)

# Slide 10: Study Comparison
add_content_slide(prs, "Major Study Results Summary", [
    "NHS (Observational): ↓ CAD risk with HRT",
    "HERS (Secondary Prevention): No difference overall",
    "  • Increased events early, protective later",
    "WHI (Primary Prevention): ↑ CVD, stroke, and VTE",
    "  • 0.07% increase overall",
    "  • Trend toward protection in 50-59 age group",
    "Danish Osteoporosis Study: Favorable outcomes when started in recently menopausal women (mean age 49.7)"
])

# Slide 11: Recommendations
add_content_slide(prs, "Clinical Recommendations", [
    "HRT is NOT recommended solely for CVD prevention",
    "HRT IS the most effective treatment for menopausal symptoms",
    "If symptoms are lifestyle-limiting:",
    "  • Start during menopausal transition",
    "  • Use lowest effective hormone doses",
    "  • Consider topical estrogens first",
    "Initiating after age 60 or continuing past 60:",
    "  • Must be individualized",
    "  • Conservative risk-benefit assessment required",
    "  • Include all standard cardiovascular risks"
])

# Slide 12: The Timing Hypothesis Summary
add_content_slide(prs, "The Timing Hypothesis", [
    "Early initiation (within menopausal transition):",
    "  • May reduce atherosclerosis progression",
    "  • Potentially cardioprotective",
    "  • Lower cardiovascular risk",
    "Late initiation (>10 years post-menopause):",
    "  • No cardiovascular benefit",
    "  • May increase risk",
    "  • Not recommended for CVD prevention",
    "Key Message: The timing of HRT initiation matters more than previously recognized"
])

# Slide 13: Conclusions
add_content_slide(prs, "Conclusions", [
    "CVD remains the #1 cause of death in US postmenopausal women",
    "The cardioprotective effect of estrogen is complex and timing-dependent",
    "Current evidence does not support HRT solely for CVD prevention",
    "Statistically significant data supporting HRT within specific age windows is accumulating",
    "For symptom management: HRT remains most effective when started early",
    "Individualized risk-benefit assessment is essential",
    "Research continues to define optimal treatment windows"
])

# Save presentation
prs.save('/Users/chileshe/Desktop/ResearchPapers/Menopause_and_the_Heart.pptx')
print("PowerPoint presentation created successfully!")
print("Saved as: Menopause_and_the_Heart.pptx")
