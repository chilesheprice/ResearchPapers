#!/usr/bin/env python3
"""
Create a PowerPoint presentation for "Menopause and the Heart: 10 Years Later" (2026 Update)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme - updated with more modern colors
TITLE_COLOR = RGBColor(13, 71, 161)  # Deep blue
SUBTITLE_COLOR = RGBColor(25, 118, 210)  # Medium blue
ACCENT_COLOR = RGBColor(198, 40, 40)  # Red for important points
TEXT_COLOR = RGBColor(54, 54, 54)  # Dark gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = SUBTITLE_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points, highlight_indices=[]):
    """Add a content slide with bullet points, optionally highlighting certain points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        # Highlight specific points in red
        if i in highlight_indices:
            p.font.color.rgb = ACCENT_COLOR
            p.font.bold = True
        else:
            p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title="", right_title=""):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    if left_title:
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_after = Pt(12)

    for i, point in enumerate(left_content):
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    if right_title:
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = SUBTITLE_COLOR
        p.space_after = Pt(12)

    for i, point in enumerate(right_content):
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

    return slide

# Slide 1: Title
add_title_slide(prs, "Menopause and the Heart: 10 Years Later",
                "A 2026 Update • Chileshe Nkonde-Price, MD")

# Slide 2: What's Changed in 10 Years
add_content_slide(prs, "What Has Changed Since 2015?", [
    "Timing hypothesis validated through multiple lines of evidence",
    "Long-term follow-up data from KEEPS, WHI, ELITE, and other major trials",
    "American Heart Association 2020 Scientific Statement on menopause and CVD",
    "Development of risk-stratified, individualized approaches to MHT",
    "FDA removes black box warnings from hormone therapy (November 2025)",
    "Shift from fear-based avoidance to evidence-based personalized care"
], highlight_indices=[4])

# Slide 3: The Timing Hypothesis - Validated
add_content_slide(prs, "The Timing Hypothesis: Now Evidence-Based", [
    "2015: Speculative concept to explain study discordances",
    "2026: Robust clinical framework supported by extensive data",
    "Key principle: Same therapy can be beneficial, neutral, or harmful depending on TIMING",
    "Early initiation (<10 years post-menopause, age <60):",
    "  • Reduced coronary heart disease (RR 0.52)",
    "  • Decreased all-cause mortality (RR 0.70)",
    "Late initiation (>10 years, age >60):",
    "  • Neutral to increased cardiovascular risk",
    "  • No benefit for primary prevention"
])

# Slide 4: Long-Term Follow-Up Data
add_content_slide(prs, "Long-Term Trial Results (2015-2026)", [
    "KEEPS Continuation (14-year follow-up):",
    "  • No increase in CVD events with early MHT",
    "  • Trend toward lower diabetes incidence",
    "  • Confirms safety of short-term early intervention",
    "Extended WHI Analysis (20+ years):",
    "  • Age 50-59 at initiation: ↓ mortality, ↓ CHD",
    "  • Age 70-79 at initiation: ↑ cardiovascular events",
    "  • Reinforces critical importance of timing",
    "Meta-analysis (40,000+ women, 2022):",
    "  • Early MHT: 48% ↓ in CHD, 30% ↓ in mortality"
])

# Slide 5: AHA 2020 Scientific Statement
add_content_slide(prs, "American Heart Association Statement (2020)", [
    "Published November 2020: landmark recognition of menopause-CVD link",
    "Key Finding: Menopausal transition itself accelerates CVD risk",
    "Vasomotor symptoms are not just quality-of-life issues:",
    "  • Associated with subclinical atherosclerosis",
    "  • Linked to adverse cardiovascular risk profiles",
    "Menopause transition = critical window for prevention",
    "Earlier age at menopause = higher CVD risk",
    "Reframed menopause as cardiovascular risk factor, not just hormone deficiency"
])

# Slide 6: FDA Black Box Warning Removal
add_content_slide(prs, "FDA Historic Decision: November 2025", [
    "BLACK BOX WARNINGS REMOVED from all MHT products",
    "Expert panel convened July 2025, decision announced November 10, 2025",
    "What's being removed:",
    "  • Warnings for cardiovascular disease",
    "  • Warnings for stroke and dementia",
    "  • Broad contraindications based on WHI data",
    "What's RETAINED:",
    "  • Endometrial cancer warning (unopposed estrogen)",
    "Implementation: New labels within 6 months (Spring 2026)"
], highlight_indices=[0])

# Slide 7: FDA Rationale
add_content_slide(prs, "Why Did the FDA Remove the Warnings?", [
    "Original warnings based on WHI population (mean age 63)",
    "Not representative of typical MHT users (ages 50-55)",
    "Age-stratified analyses showed younger women benefited or had neutral outcomes",
    "Absolute risk increases were small, even in older WHI participants",
    "Black box warnings created 'chilling effect':",
    "  • 66% decline in prescriptions post-2002",
    "  • Millions of women suffered with untreated symptoms",
    "Modern evidence supports safety in appropriately selected women",
    "Labels now include age-specific guidance emphasizing timing"
])

# Slide 8: Risk Stratification Framework
add_content_slide(prs, "Contemporary Risk Stratification (2026)", [
    "LOW RISK: Age <60, <10 years post-menopause, no CVD/risk factors, CAC=0",
    "  → MHT appropriate for symptoms",
    "INTERMEDIATE RISK: Age <60, 1-2 risk factors, CAC 1-99, ASCVD <10%",
    "  → MHT reasonable with shared decision-making",
    "HIGH RISK: Age ≥60, multiple risk factors, CAC 100-400, ASCVD 10-20%",
    "  → MHT generally not recommended; consider alternatives",
    "VERY HIGH RISK: Established ASCVD, CAC >400, ASCVD >20%, VTE history",
    "  → MHT CONTRAINDICATED regardless of symptoms"
])

# Slide 9: Route and Formulation Matters
add_two_column_slide(prs, "Route of Administration: Critical Consideration",
    [
        "TRANSDERMAL ESTROGEN",
        "Preferred in 2026 guidelines",
        "Benefits:",
        "• 50% ↓ VTE risk vs oral",
        "• Lower stroke risk",
        "• Less effect on clotting factors",
        "• No first-pass hepatic metabolism",
        "• Better lipid profile",
        "Ideal for:",
        "• All women if tolerated",
        "• Especially: obesity, age >60,",
        "  diabetes, migraine, high TG"
    ],
    [
        "ORAL ESTROGEN",
        "Alternative if transdermal not tolerated",
        "Concerns:",
        "• Higher VTE risk",
        "• Increases triglycerides",
        "• ↑ inflammatory markers (CRP)",
        "• First-pass liver effects",
        "",
        "Progestogen:",
        "• Micronized progesterone preferred",
        "• More neutral CV profile vs synthetic",
        "• 200 mg nightly for intact uterus"
    ],
    "TRANSDERMAL", "ORAL"
)

# Slide 10: 2026 Recommendations - Symptomatic Women
add_content_slide(prs, "2026 Recommendations: Early Menopausal Women", [
    "FOR WOMEN <60 YEARS OR <10 YEARS POST-MENOPAUSE:",
    "MHT is APPROPRIATE for moderate-to-severe vasomotor symptoms",
    "Do NOT withhold due to cardiovascular concerns in low-risk women",
    "Preferred regimen:",
    "  • Transdermal estradiol 0.025-0.05 mg/day (lowest effective dose)",
    "  • + Micronized progesterone 200 mg nightly (if uterus intact)",
    "Duration: Continue as long as benefits outweigh risks",
    "  • No arbitrary time limit",
    "  • Re-evaluate annually",
    "  • Consider tapering after 5-7 years if asymptomatic"
])

# Slide 11: 2026 Recommendations - Older Women
add_content_slide(prs, "2026 Recommendations: Late Menopausal Women", [
    "FOR WOMEN >60 YEARS OR >10 YEARS POST-MENOPAUSE:",
    "Do NOT initiate MHT without comprehensive cardiovascular assessment",
    "If severe persistent symptoms:",
    "  • Perform thorough CV risk stratification",
    "  • Consider CAC scoring if available",
    "  • If low-intermediate risk: shared decision-making may support trial",
    "  • Use lowest dose transdermal estrogen (0.025 mg/day)",
    "If already on MHT and well-tolerated:",
    "  • Continuation reasonable in low-risk women",
    "  • Reassess CV risk annually",
    "Consider non-hormonal alternatives first"
])

# Slide 12: Absolute Contraindications
add_content_slide(prs, "Absolute Contraindications to MHT (2026)", [
    "Established cardiovascular disease:",
    "  • Prior MI, stroke, TIA, or revascularization",
    "Active or history of venous thromboembolism (VTE)",
    "Known thrombophilia",
    "Active liver disease",
    "Estrogen-dependent malignancies (current or recent)",
    "Undiagnosed vaginal bleeding",
    "Pregnancy",
    "",
    "NOTE: Removal of FDA black box warnings does NOT mean MHT is safe for everyone"
])

# Slide 13: Special Population - Premature Menopause
add_content_slide(prs, "Special Case: Premature/Early Menopause", [
    "Age <40 (premature) or <45 (early menopause):",
    "MHT is STRONGLY RECOMMENDED until at least age 50",
    "Rationale:",
    "  • Mitigates long-term cardiovascular risk",
    "  • Prevents accelerated bone loss",
    "  • Protects cognitive function",
    "  • Reduces all-cause mortality",
    "Higher doses appropriate (replicate premenopausal levels)",
    "Standard cardiovascular contraindications do NOT apply",
    "This is hormone replacement, not supplementation"
])

# Slide 14: Clinical Case 1
add_content_slide(prs, "Case 1: Recently Menopausal, Severe Symptoms", [
    "52-year-old, 18 months post-menopause",
    "Severe hot flashes (15-20/day), night sweats, sleep disruption",
    "No medical history, normal BMI, BP 118/72, LDL 105, glucose 92",
    "Quality of life severely impacted",
    "",
    "RECOMMENDATION:",
    "✓ IDEAL candidate for MHT",
    "✓ Transdermal estradiol 0.05 mg/day",
    "✓ Micronized progesterone 200 mg nightly",
    "✓ Minimal CV risk at her age and timing",
    "✓ Continue as long as benefits persist, no arbitrary time limit"
])

# Slide 15: Clinical Case 2
add_content_slide(prs, "Case 2: Remote from Menopause, Wants MHT Now", [
    "64-year-old, 16 years post-menopause",
    "Persistent hot flashes, avoided MHT due to old black box warnings",
    "Well-controlled hypertension (125/78), BMI 28, LDL 115, no diabetes",
    "Requests MHT now that warnings removed",
    "",
    "RECOMMENDATION:",
    "⚠ Requires careful individualized assessment",
    "• Calculate 10-year ASCVD risk (likely 5-10% = low-intermediate)",
    "• Obtain CAC score if available",
    "• If CAC <100: shared decision-making may support trial",
    "• Low-dose transdermal estradiol 0.025 mg/day",
    "• Also discuss non-hormonal alternatives (fezolinetant, paroxetine)",
    "• Aggressive risk factor management essential"
])

# Slide 16: Clinical Case 3
add_content_slide(prs, "Case 3: Established CVD, Severe Symptoms", [
    "58-year-old, 6 years post-menopause",
    "Non-STEMI 2 years ago, stents in place",
    "On dual antiplatelet therapy, statin, beta-blocker",
    "Severe vasomotor symptoms impacting quality of life",
    "Asks about MHT now that 'warnings are gone'",
    "",
    "RECOMMENDATION:",
    "✗ MHT remains CONTRAINDICATED in established CVD",
    "✗ FDA label change does NOT make MHT safe for everyone",
    "✓ Recommend non-hormonal approaches:",
    "   • Fezolinetant (neurokinin-3 receptor antagonist)",
    "   • Paroxetine, gabapentin, CBT, lifestyle modifications",
    "✓ Optimize cardiovascular medications and risk factors"
])

# Slide 17: The Paradigm Shift
add_content_slide(prs, "The Paradigm Shift: 2002 to 2026", [
    "2002-2015: Fear-based universal avoidance",
    "  • 'Hormone therapy causes heart disease'",
    "  • MHT prescriptions dropped 66%",
    "  • Millions suffered unnecessarily",
    "",
    "2026: Evidence-based individualized care",
    "  • 'Timing is paramount'",
    "  • MHT appropriate for right patient at right time",
    "  • Comprehensive risk stratification",
    "  • Route and formulation matter",
    "  • Shared decision-making essential",
    "",
    "From 'one-size-fits-all prohibition' to 'personalized medicine'"
])

# Slide 18: Key Principles for 2026
add_content_slide(prs, "Five Key Principles (2026)", [
    "1. TIMING IS PARAMOUNT",
    "   Same intervention → different outcomes based on when initiated",
    "",
    "2. MHT IS NOT CARDIOVASCULAR PREVENTION",
    "   Don't prescribe solely to prevent heart disease",
    "",
    "3. ROUTE MATTERS",
    "   Transdermal safer than oral for most CV outcomes",
    "",
    "4. INDIVIDUALIZATION IS ESSENTIAL",
    "   Risk assessment + shared decision-making required",
    "",
    "5. SYMPTOMS MATTER",
    "   Severe VMS associated with CV risk and justify treatment"
])

# Slide 19: Ongoing Questions
add_content_slide(prs, "Remaining Questions and Future Research", [
    "Optimal duration: Do long-term users (15-20 years) maintain benefits?",
    "Personalized medicine: Can genetics identify who will benefit most?",
    "Contemporary formulations: Do transdermal + micronized progesterone show superior outcomes in RCTs?",
    "Heart failure: What is MHT's impact on HFpEF in postmenopausal women?",
    "Racial/ethnic differences: Does timing hypothesis apply equally across diverse populations?",
    "Non-hormonal alternatives: Long-term CV safety of fezolinetant and other new agents?"
])

# Slide 20: Conclusions
add_content_slide(prs, "Conclusions: 10 Years of Progress", [
    "The timing hypothesis has been validated and refined",
    "FDA's November 2025 black box warning removal marks new era in women's health",
    "For appropriate candidates (age <60, <10 years post-menopause):",
    "  • MHT is safe and effective",
    "  • Does not increase—may reduce—cardiovascular risk",
    "For women remote from menopause or with CVD:",
    "  • MHT generally not appropriate",
    "The era of one-size-fits-all recommendations is over",
    "The future is individualized, evidence-based, woman-centered care",
    "",
    "That future has arrived."
])

# Save presentation
prs.save('/Users/chileshe/Desktop/ResearchPapers/Menopause_and_the_Heart_2026_Update.pptx')
print("PowerPoint presentation created successfully!")
print("Saved as: Menopause_and_the_Heart_2026_Update.pptx")
